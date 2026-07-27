import io
import math
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class ExportacaoPPTXService:
    """
    Serviço avançado de geração de apresentações PowerPoint (.pptx).
    Aplica design corporativo idêntico ao Tailwind e suporta paginação dinâmica de tabelas.
    """

    COLOR_BG = RGBColor(248, 250, 252)       # Slate 50 (#F8FAFC)
    COLOR_PRIMARY = RGBColor(11, 37, 69)     # Dark Blue Navy (#0B2545)
    COLOR_TEXT = RGBColor(15, 23, 42)        # Slate 900 (#0F172A)
    COLOR_MUTED = RGBColor(100, 116, 139)    # Slate 500 (#64748B)
    COLOR_CARD_BG = RGBColor(255, 255, 255)  # White (#FFFFFF)
    COLOR_ACCENT = RGBColor(37, 99, 235)     # Blue 600 (#2563EB)
    COLOR_BORDER = RGBColor(226, 232, 240)   # Slate 200 (#E2E8F0)

    @staticmethod
    def _definir_fundo_slide(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = ExportacaoPPTXService.COLOR_BG

    @staticmethod
    def _adicionar_cabecalho(slide, titulo_texto, subtitulo_texto=""):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = titulo_texto.upper()
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ExportacaoPPTXService.COLOR_PRIMARY
        p.font.name = "Arial"

        if subtitulo_texto:
            p2 = tf.add_paragraph()
            p2.text = subtitulo_texto
            p2.font.size = Pt(12)
            p2.font.color.rgb = ExportacaoPPTXService.COLOR_MUTED
            p2.font.name = "Arial"

    @staticmethod
    def gerar_apresentacao_completa(emissor_nome, kpis, lancamentos, status_counts, setor_nome="Todos os Setores"):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # -------------------------------------------------------------
        # SLIDE 1: CAPA
        # -------------------------------------------------------------
        slide_capa = prs.slides.add_slide(blank_layout)
        ExportacaoPPTXService._definir_fundo_slide(slide_capa)

        # Barra lateral decorativa
        shape_bar = slide_capa.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
        )
        shape_bar.fill.solid()
        shape_bar.fill.fore_color.rgb = ExportacaoPPTXService.COLOR_PRIMARY
        shape_bar.line.color.rgb = ExportacaoPPTXService.COLOR_PRIMARY

        # Caixa do Título Principal
        tb_title = slide_capa.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.5))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True

        p1 = tf_title.paragraphs[0]
        p1.text = "Relatório de Gestão Operacional"
        p1.font.size = Pt(38)
        p1.font.bold = True
        p1.font.color.rgb = ExportacaoPPTXService.COLOR_PRIMARY
        p1.font.name = "Arial"

        p2 = tf_title.add_paragraph()
        p2.text = f"Consolidado de Indicadores, SLAs e Desempenho — {setor_nome}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = ExportacaoPPTXService.COLOR_MUTED
        p2.font.name = "Arial"

        # Rodapé da Capa
        tb_meta = slide_capa.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(11.0), Inches(1.2))
        tf_meta = tb_meta.text_frame
        p3 = tf_meta.paragraphs[0]
        p3.text = f"Gerado por: {emissor_nome} | Data: {datetime.now().strftime('%d/%m/%Y às %H:%M')}"
        p3.font.size = Pt(12)
        p3.font.bold = True
        p3.font.color.rgb = ExportacaoPPTXService.COLOR_TEXT
        p3.font.name = "Arial"

        # -------------------------------------------------------------
        # SLIDE 2: INDICADORES E PAINEL DE STATUS
        # -------------------------------------------------------------
        slide_kpis = prs.slides.add_slide(blank_layout)
        ExportacaoPPTXService._definir_fundo_slide(slide_kpis)
        ExportacaoPPTXService._adicionar_cabecalho(slide_kpis, "Visão Geral de Indicadores", "Métricas de produção e atratividade do período")

        card_configs = [
            ("HORAS PRODUZIDAS", kpis.get('horas', '0h'), Inches(0.8), Inches(1.8)),
            ("EFICIÊNCIA MÉDIA", kpis.get('eficiencia', '0%'), Inches(3.8), Inches(1.8)),
            ("COLABORADORES", kpis.get('usuarios', '0'), Inches(6.8), Inches(1.8)),
            ("ROTINAS ATIVAS", kpis.get('rotinas', '0'), Inches(9.8), Inches(1.8))
        ]

        for label, val, left_pos, top_pos in card_configs:
            card = slide_kpis.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, Inches(2.7), Inches(1.5))
            card.fill.solid()
            card.fill.fore_color.rgb = ExportacaoPPTXService.COLOR_CARD_BG
            card.line.color.rgb = ExportacaoPPTXService.COLOR_BORDER

            tf_c = card.text_frame
            tf_c.word_wrap = True
            p_lbl = tf_c.paragraphs[0]
            p_lbl.text = label
            p_lbl.font.size = Pt(10)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = ExportacaoPPTXService.COLOR_MUTED

            p_val = tf_c.add_paragraph()
            p_val.text = str(val)
            p_val.font.size = Pt(24)
            p_val.font.bold = True
            p_val.font.color.rgb = ExportacaoPPTXService.COLOR_PRIMARY

        # Painel SLAs
        tb_sla_box = slide_kpis.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.8), Inches(11.7), Inches(2.8))
        tb_sla_box.fill.solid()
        tb_sla_box.fill.fore_color.rgb = ExportacaoPPTXService.COLOR_CARD_BG
        tb_sla_box.line.color.rgb = ExportacaoPPTXService.COLOR_BORDER

        tf_sla = tb_sla_box.text_frame
        p_sla_title = tf_sla.paragraphs[0]
        p_sla_title.text = "DISTRIBUIÇÃO DE ATIVIDADES POR STATUS SLA"
        p_sla_title.font.size = Pt(14)
        p_sla_title.font.bold = True
        p_sla_title.font.color.rgb = ExportacaoPPTXService.COLOR_PRIMARY

        status_items = [
            ("Concluídos", status_counts.get('Concluído', 0)),
            ("Em Andamento", status_counts.get('Em Andamento', 0)),
            ("Cancelados", status_counts.get('Cancelado', 0))
        ]

        for idx, (st_label, st_val) in enumerate(status_items):
            p_st = tf_sla.add_paragraph()
            p_st.text = f"• {st_label}: {st_val} atividade(s)"
            p_st.font.size = Pt(14)
            p_st.font.color.rgb = ExportacaoPPTXService.COLOR_TEXT

        # -------------------------------------------------------------
        # SLIDE(S) 3+: TABELAS DE DETALHAMENTO COM PAGINAÇÃO DINÂMICA
        # -------------------------------------------------------------
        ITENS_POR_SLIDE = 8
        total_registros = len(lancamentos)
        total_paginas = max(1, math.ceil(total_registros / ITENS_POR_SLIDE))

        for pag in range(total_paginas):
            slide_tb = prs.slides.add_slide(blank_layout)
            ExportacaoPPTXService._definir_fundo_slide(slide_tb)
            ExportacaoPPTXService._adicionar_cabecalho(
                slide_tb, 
                "Detalhamento de Lançamentos", 
                f"Página {pag + 1} de {total_paginas} — Exibindo registros de produção"
            )

            inicio_idx = pag * ITENS_POR_SLIDE
            fim_idx = inicio_idx + ITENS_POR_SLIDE
            lote_lancamentos = lancamentos[inicio_idx:fim_idx]

            num_rows = len(lote_lancamentos) + 1
            num_cols = 6
            table_shape = slide_tb.shapes.add_table(
                num_rows, num_cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.5 * num_rows)
            )
            table = table_shape.table

            headers = ["Data", "Colaborador", "Setor", "Atividade", "Duração", "Eficiência"]
            col_widths = [Inches(1.5), Inches(3.0), Inches(2.0), Inches(3.2), Inches(1.0), Inches(1.0)]

            for col_idx, width in enumerate(col_widths):
                table.columns[col_idx].width = width

            for col_idx, h_text in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = ExportacaoPPTXService.COLOR_PRIMARY
                p = cell.text_frame.paragraphs[0]
                p.text = h_text
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

            for row_idx, lanc in enumerate(lote_lancamentos, start=1):
                dt_str = lanc.data_programada.strftime('%d/%m/%Y') if lanc.data_programada else ''
                autor_str = lanc.autor.nome_completo if lanc.autor else 'N/A'
                setor_str = lanc.setor_snapshot.sigla if lanc.setor_snapshot else 'N/A'
                atv_str = lanc.atividade_referencia.titulo if lanc.atividade_referencia else 'N/A'
                dur_str = f"{lanc.duracao_minutos} min"
                efic_str = f"{round(lanc.eficiencia_percentual or 0, 1)}%"

                row_values = [dt_str, autor_str, setor_str, atv_str, dur_str, efic_str]

                for col_idx, val_text in enumerate(row_values):
                    cell = table.cell(row_idx, col_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ExportacaoPPTXService.COLOR_CARD_BG if row_idx % 2 == 0 else ExportacaoPPTXService.COLOR_BG
                    p = cell.text_frame.paragraphs[0]
                    p.text = val_text
                    p.font.size = Pt(10)
                    p.font.color.rgb = ExportacaoPPTXService.COLOR_TEXT
                    if col_idx in [0, 2, 4, 5]:
                        p.alignment = PP_ALIGN.CENTER

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output